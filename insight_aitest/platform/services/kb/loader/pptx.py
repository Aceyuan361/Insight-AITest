# -*- coding: utf-8 -*-
"""PowerPoint 文档 Loader（.pptx）。用 python-pptx 提取每张幻灯片的文本与表格。"""

from __future__ import annotations

from pathlib import Path

from insight_aitest.platform.services.kb.loader import register_loader
from insight_aitest.platform.services.kb.loader.base import DocumentLoader
from insight_aitest.platform.services.kb.models import ParsedDocument


@register_loader(".pptx")
class PptxLoader(DocumentLoader):
    """PowerPoint .pptx 文档 Loader。逐幻灯片提取标题、正文文本框与表格。

    注意：全量加载到内存（python-pptx 无流式模式），依赖上传端点的
    max_upload_mb 检查限制单文件大小（默认 20MB）。
    """

    def load(self, path: Path) -> ParsedDocument:
        from pptx import Presentation

        prs = Presentation(str(path))
        parts: list[str] = []

        for idx, slide in enumerate(prs.slides, start=1):
            slide_texts: list[str] = []

            # 标题（若有）
            if slide.shapes.title is not None and slide.shapes.title.text:
                slide_texts.append(slide.shapes.title.text.strip())

            # 文本框 / 占位符
            for shape in slide.shapes:
                # 标题已在上面处理，跳过
                if shape.has_text_frame and shape != slide.shapes.title:
                    text = shape.text_frame.text.strip()
                    if text:
                        slide_texts.append(text)
                # 表格
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            slide_texts.append(" | ".join(cells))

            if slide_texts:
                parts.append(f"--- 第 {idx} 页 ---")
                parts.extend(slide_texts)

        content = "\n".join(parts)
        return ParsedDocument(filename=path.name, content=content, meta={"format": "pptx"})
