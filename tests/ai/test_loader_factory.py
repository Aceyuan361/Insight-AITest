# -*- coding: utf-8 -*-
import pytest

from insight_aitest.platform.services.kb.loader import get_loader
from insight_aitest.platform.services.kb.loader.base import (
    DocumentLoader, UnsupportedFormatError,
)


def test_get_loader_md():
    loader = get_loader("readme.md")
    assert isinstance(loader, DocumentLoader)


def test_get_loader_txt():
    assert isinstance(get_loader("a.txt"), DocumentLoader)


def test_get_loader_pdf():
    assert isinstance(get_loader("a.pdf"), DocumentLoader)


def test_get_loader_docx():
    """docx 已注册（python-docx 解析）。"""
    assert isinstance(get_loader("a.docx"), DocumentLoader)


def test_get_loader_unknown_raises():
    with pytest.raises(UnsupportedFormatError):
        get_loader("a.xyz")


def test_get_loader_image_extensions():
    """图片扩展名已注册（OCR 走 vision model）。"""
    for ext in (".png", ".jpg", ".jpeg", ".gif", ".webp"):
        assert isinstance(get_loader(f"a{ext}"), DocumentLoader), f"{ext} 未注册"


def test_get_loader_image_case_insensitive():
    assert isinstance(get_loader("PHOTO.PNG"), DocumentLoader)
    assert isinstance(get_loader("PHOTO.JPG"), DocumentLoader)


def test_extension_case_insensitive():
    assert isinstance(get_loader("A.PDF"), DocumentLoader)


from pathlib import Path  # noqa: E402

FIXTURES = Path(__file__).parent / "fixtures"


def test_markdown_loader():
    from insight_aitest.platform.services.kb.loader import get_loader
    doc = get_loader("sample.md").load(FIXTURES / "sample.md")
    assert "# 测试文档" in doc.content
    assert "这是第一段内容" in doc.content
    assert doc.filename == "sample.md"
    assert doc.meta.get("char_count", 0) > 0


def test_text_loader_utf8():
    from insight_aitest.platform.services.kb.loader import get_loader
    doc = get_loader("sample.txt").load(FIXTURES / "sample.txt")
    assert "这是一段纯文本" in doc.content
    assert doc.meta.get("char_count", 0) > 0


# ===== Office / HTML 格式注册（无需第三方库，始终跑）=====


def test_get_loader_xlsx():
    assert isinstance(get_loader("a.xlsx"), DocumentLoader)


def test_get_loader_pptx():
    assert isinstance(get_loader("a.pptx"), DocumentLoader)


def test_get_loader_html():
    assert isinstance(get_loader("a.html"), DocumentLoader)


def test_get_loader_htm_alias():
    """htm 别名注册同一个 Loader。"""
    assert isinstance(get_loader("a.htm"), DocumentLoader)


def test_get_loader_office_case_insensitive():
    assert isinstance(get_loader("REPORT.XLSX"), DocumentLoader)
    assert isinstance(get_loader("DECK.PPTX"), DocumentLoader)
    assert isinstance(get_loader("PAGE.HTML"), DocumentLoader)


# ===== 功能测试（依赖第三方库，未装则优雅跳过）=====


def test_html_loader_extracts_text(tmp_path):
    """HTML loader：提取可见文本，剔除 script/style。"""
    pytest.importorskip("bs4")
    from insight_aitest.platform.services.kb.loader import get_loader

    html = """<!DOCTYPE html>
<html><head><title>测试页</title>
<style>body{color:red}</style>
<script>alert('x')</script>
</head><body>
<h1>需求规格说明书</h1>
<p>支持账号密码和手机验证码登录。</p>
<div>隐藏的<script>evil()</script>内容</div>
</body></html>"""
    p = tmp_path / "spec.html"
    p.write_text(html, encoding="utf-8")

    doc = get_loader("spec.html").load(p)
    assert "需求规格说明书" in doc.content
    assert "账号密码" in doc.content
    assert doc.meta["format"] == "html"
    # script/style 内容应被剔除
    assert "alert" not in doc.content
    assert "evil" not in doc.content
    assert "color:red" not in doc.content


def test_xlsx_loader_extracts_rows(tmp_path):
    """Excel loader：逐行提取单元格，用 ` | ` 拼接。"""
    pytest.importorskip("openpyxl")
    from openpyxl import Workbook
    from insight_aitest.platform.services.kb.loader import get_loader

    wb = Workbook()
    ws = wb.active
    ws.title = "登录用例"
    ws.append(["用例名", "步骤", "预期"])
    ws.append(["正常登录", "输入正确账号", "登录成功"])
    ws.append(["密码错误", "输入错误密码", "提示密码错误"])
    p = tmp_path / "cases.xlsx"
    wb.save(p)

    doc = get_loader("cases.xlsx").load(p)
    assert "登录用例" in doc.content  # sheet 名
    assert "用例名 | 步骤 | 预期" in doc.content  # 表头行
    assert "正常登录 | 输入正确账号 | 登录成功" in doc.content
    assert "密码错误" in doc.content
    assert doc.meta["format"] == "xlsx"


def test_pptx_loader_extracts_slides(tmp_path):
    """PowerPoint loader：逐页提取标题 + 正文 + 表格。"""
    pytest.importorskip("pptx")
    from pptx import Presentation
    from insight_aitest.platform.services.kb.loader import get_loader

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # 标题版式
    slide.shapes.title.text = "测试概览"
    slide.placeholders[1].text = "本次覆盖登录与支付模块"
    p = tmp_path / "overview.pptx"
    prs.save(p)

    doc = get_loader("overview.pptx").load(p)
    assert "测试概览" in doc.content
    assert "登录与支付模块" in doc.content
    assert doc.meta["format"] == "pptx"
